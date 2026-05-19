from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import os

def px(n): return int(n * 9525)

# Exact colors from fragrance_notes_template.pptx
CREAM    = RGBColor(0xF5, 0xF2, 0xEC)
DIVIDER  = RGBColor(0xC4, 0xB9, 0xA8)
GOLD     = RGBColor(0xBF, 0x9A, 0x58)
BROWN    = RGBColor(0x8B, 0x73, 0x55)
PHOTO_BG = RGBColor(0xE8, 0xE2, 0xD8)

# Slide: 768 × 960 px (matches template)
SLIDE_W, SLIDE_H = 768, 960

# 2-per-row layout (from TOP NOTES section in template)
PHOTO_W = 125
PHOTO_X = [240, 403]   # x positions for 2 photos per row

# Section layout (y positions + header/divider x positions per section)
# Taken directly from fragrance_notes_template.pptx measurements
SECTIONS = {
    "TOP":    {"header": 75,  "divider": 86,  "photo": 127, "label": 267,
               "div_l_x": 154, "div_l_w": 154, "hdr_x": 302, "hdr_w": 163, "div_r_x": 461, "div_r_w": 154},
    "MIDDLE": {"header": 372, "divider": 384, "photo": 424, "label": 564,
               "div_l_x": 134, "div_l_w": 154, "hdr_x": 298, "hdr_w": 173, "div_r_x": 480, "div_r_w": 154},
    "BASE":   {"header": 670, "divider": 682, "photo": 722, "label": 862,
               "div_l_x": 144, "div_l_w": 154, "hdr_x": 302, "hdr_w": 163, "div_r_x": 470, "div_r_w": 154},
}


def build(fragrance_name, sections, output_path):
    prs = Presentation()
    prs.slide_width  = px(SLIDE_W)
    prs.slide_height = px(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg_shape = slide.shapes.add_shape(1, 0, 0, px(SLIDE_W), px(SLIDE_H))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = CREAM
    bg_shape.line.fill.background()

    for key, notes in sections.items():
        y = SECTIONS[key]

        # Left divider
        ln = slide.shapes.add_shape(1, px(y["div_l_x"]), px(y["divider"]), px(y["div_l_w"]), px(2))
        ln.fill.solid()
        ln.fill.fore_color.rgb = DIVIDER
        ln.line.fill.background()

        # Right divider
        rn = slide.shapes.add_shape(1, px(y["div_r_x"]), px(y["divider"]), px(y["div_r_w"]), px(2))
        rn.fill.solid()
        rn.fill.fore_color.rgb = DIVIDER
        rn.line.fill.background()

        # Header text
        tb = slide.shapes.add_textbox(px(y["hdr_x"]), px(y["header"]), px(y["hdr_w"]), px(25))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = key + " NOTES"
        rPr = run._r.get_or_add_rPr()
        rPr.set("sz", "1100")
        rPr.set("b", "0")
        rPr.set("spc", "200")
        fill_el = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        clr_el  = etree.SubElement(fill_el, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
        clr_el.set("val", "BF9A58")
        lat_el = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
        lat_el.set("typeface", "Montserrat")

        # Photo placeholders + labels
        for i, (ingredient, _) in enumerate(notes):
            x = PHOTO_X[i]

            # Rounded rect placeholder
            sp = slide.shapes.add_shape(5, px(x), px(y["photo"]), px(PHOTO_W), px(PHOTO_W))
            sp.fill.solid()
            sp.fill.fore_color.rgb = PHOTO_BG
            sp.line.fill.background()
            # Set corner radius to match template (adj=8000)
            prstGeom = sp._element.spPr.prstGeom
            avLst = prstGeom.find("{http://schemas.openxmlformats.org/drawingml/2006/main}avLst")
            if avLst is None:
                avLst = etree.SubElement(prstGeom, "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst")
            gd = etree.SubElement(avLst, "{http://schemas.openxmlformats.org/drawingml/2006/main}gd")
            gd.set("name", "adj")
            gd.set("fmla", "val 8000")

            # Label
            lb = slide.shapes.add_textbox(px(x), px(y["label"]), px(PHOTO_W), px(23))
            ltf = lb.text_frame
            ltf.word_wrap = False
            lp = ltf.paragraphs[0]
            lp.alignment = PP_ALIGN.CENTER
            lrun = lp.add_run()
            lrun.text = ingredient
            lrPr = lrun._r.get_or_add_rPr()
            lrPr.set("sz", "1400")
            lrPr.set("b", "0")
            lfill = etree.SubElement(lrPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
            lclr  = etree.SubElement(lfill, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
            lclr.set("val", "8B7355")
            llat = etree.SubElement(lrPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
            llat.set("typeface", "Montserrat")

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    build(
        fragrance_name="Lift Me Up",
        sections={
            "TOP":    [("Magnolia", None), ("Bergamot", None)],
            "MIDDLE": [("Tonka", None),    ("Ylang Ylang", None)],
            "BASE":   [("Vanilla", None),  ("Musk", None)],
        },
        output_path=os.path.expanduser("~/Desktop/initio-lift-me-up.pptx")
    )
